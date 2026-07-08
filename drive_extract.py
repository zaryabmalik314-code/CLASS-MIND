"""
ClassMind - Drive file text extraction.

Given a Drive file id, figures out its type and pulls plain text out of
it. Handles the three formats most likely for LGU slide uploads: native
Google Slides, PDF, and PPTX. Anything else raises ValueError so the
poller can log and skip it instead of crashing the whole loop.
"""

import io

from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

GOOGLE_SLIDES_MIME = "application/vnd.google-apps.presentation"
PDF_MIME = "application/pdf"
PPTX_MIME = (
    "application/vnd.openxmlformats-officedocument.presentationml.presentation"
)


def get_drive_service(creds):
    return build("drive", "v3", credentials=creds)


def _download_bytes(drive_service, file_id: str) -> bytes:
    request = drive_service.files().get_media(fileId=file_id)
    buf = io.BytesIO()
    downloader = MediaIoBaseDownload(buf, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()
    return buf.getvalue()


def _extract_pdf_text(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    return "\n\n".join(page.extract_text() or "" for page in reader.pages)


def _extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides_text = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"--- Slide {i} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
        slides_text.append("\n".join(lines))
    return "\n\n".join(slides_text)


def extract_text(creds, file_id: str) -> tuple[str, str]:
    """Returns (title, extracted_text). Raises ValueError if unsupported."""
    drive_service = get_drive_service(creds)
    meta = (
        drive_service.files()
        .get(fileId=file_id, fields="id,name,mimeType")
        .execute()
    )
    title = meta["name"]
    mime = meta["mimeType"]

    if mime == GOOGLE_SLIDES_MIME:
        # Native Slides export straight to plain text — simplest path,
        # avoids parsing the Slides API's nested page/element JSON.
        data = drive_service.files().export(
            fileId=file_id, mimeType="text/plain"
        ).execute()
        text = data.decode("utf-8") if isinstance(data, bytes) else data
        return title, text

    if mime == PDF_MIME:
        data = _download_bytes(drive_service, file_id)
        return title, _extract_pdf_text(data)

    if mime == PPTX_MIME:
        data = _download_bytes(drive_service, file_id)
        return title, _extract_pptx_text(data)

    raise ValueError(f"Unsupported file type for '{title}': {mime}")
